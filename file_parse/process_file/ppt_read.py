# !/usr/bin/env python
# -*- coding:utf-8 -*-　
# @Time : 2023/5/13 17:41 
# @Author : sanmaomashi
# @GitHub : https://github.com/sanmaomashi

import pptx
from typing import Dict
from utils.common_util import cost_time


@cost_time
async def read_ppt_file(file_path) -> Dict:
    """
    读取指定路径下的 PowerPoint 文件内容，并返回一个文本数据列表。
    :return: 列表，其中每个元素是 PowerPoint 文件中幻灯片文本框中的文本内容。
    """
    text_list = []

    prs = pptx.Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_list.append(shape.text)
    response_data = {
        "file": file_path,
        "content": text_list,
    }
    return response_data
